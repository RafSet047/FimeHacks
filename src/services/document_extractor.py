import os
from pathlib import Path
from typing import Optional, Dict, Any
import logging
from abc import ABC, abstractmethod

# Document processing imports
try:
    from langchain_community.document_loaders import PyPDFLoader, UnstructuredWordDocumentLoader, TextLoader, CSVLoader, UnstructuredPowerPointLoader
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    LANGCHAIN_AVAILABLE = True
except ImportError:
    LANGCHAIN_AVAILABLE = False
    logging.warning("langchain not available, using fallback document processing")

# Fallback imports
try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    from docx import Document
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

try:
    from pydub import AudioSegment
    PYDUB_AVAILABLE = True
except ImportError:
    PYDUB_AVAILABLE = False

from src.config.settings import settings

logger = logging.getLogger(__name__)


class DocumentExtractor(ABC):
    """Abstract base class for document extractors"""
    
    @abstractmethod
    def extract_text(self, file_path: str) -> str:
        pass
    
    @abstractmethod
    def can_handle(self, file_path: str, mime_type: str) -> bool:
        pass


class PDFExtractor(DocumentExtractor):
    """Extract text from PDF files"""
    
    def can_handle(self, file_path: str, mime_type: str) -> bool:
        return (
            mime_type == "application/pdf" or 
            file_path.lower().endswith('.pdf')
        )
    
    def extract_text(self, file_path: str) -> str:
        if LANGCHAIN_AVAILABLE:
            return self._extract_with_langchain(file_path)
        elif PYPDF2_AVAILABLE:
            return self._extract_with_pypdf2(file_path)
        else:
            raise RuntimeError("No PDF processing library available")
    
    def _extract_with_langchain(self, file_path: str) -> str:
        try:
            loader = PyPDFLoader(file_path)
            documents = loader.load()
            return "\n\n".join([doc.page_content for doc in documents])
        except Exception as e:
            logger.error(f"Langchain PDF extraction failed: {e}")
            if PYPDF2_AVAILABLE:
                return self._extract_with_pypdf2(file_path)
            raise
    
    def _extract_with_pypdf2(self, file_path: str) -> str:
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"PyPDF2 extraction failed: {e}")
            raise


class DOCXExtractor(DocumentExtractor):
    """Extract text from DOCX files"""
    
    def can_handle(self, file_path: str, mime_type: str) -> bool:
        return (
            mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or
            file_path.lower().endswith('.docx')
        )
    
    def extract_text(self, file_path: str) -> str:
        if LANGCHAIN_AVAILABLE:
            return self._extract_with_langchain(file_path)
        elif PYTHON_DOCX_AVAILABLE:
            return self._extract_with_python_docx(file_path)
        else:
            raise RuntimeError("No DOCX processing library available")
    
    def _extract_with_langchain(self, file_path: str) -> str:
        try:
            loader = UnstructuredWordDocumentLoader(file_path)
            documents = loader.load()
            return "\n\n".join([doc.page_content for doc in documents])
        except Exception as e:
            logger.error(f"Langchain DOCX extraction failed: {e}")
            if PYTHON_DOCX_AVAILABLE:
                return self._extract_with_python_docx(file_path)
            raise
    
    def _extract_with_python_docx(self, file_path: str) -> str:
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"python-docx extraction failed: {e}")
            raise


class DOCExtractor(DocumentExtractor):
    """Extract text from DOC files"""
    
    def can_handle(self, file_path: str, mime_type: str) -> bool:
        return (
            mime_type == "application/msword" or
            file_path.lower().endswith('.doc')
        )
    
    def extract_text(self, file_path: str) -> str:
        if LANGCHAIN_AVAILABLE:
            return self._extract_with_langchain(file_path)
        else:
            raise RuntimeError("DOC file processing requires langchain with unstructured")
    
    def _extract_with_langchain(self, file_path: str) -> str:
        try:
            loader = UnstructuredWordDocumentLoader(file_path)
            documents = loader.load()
            return "\n\n".join([doc.page_content for doc in documents])
        except Exception as e:
            logger.error(f"Langchain DOC extraction failed: {e}")
            raise


class CSVExtractor(DocumentExtractor):
    """Extract text from CSV files"""
    
    def can_handle(self, file_path: str, mime_type: str) -> bool:
        return (
            mime_type == "text/csv" or
            file_path.lower().endswith('.csv')
        )
    
    def extract_text(self, file_path: str) -> str:
        if LANGCHAIN_AVAILABLE:
            return self._extract_with_langchain(file_path)
        elif PANDAS_AVAILABLE:
            return self._extract_with_pandas(file_path)
        else:
            return self._extract_with_builtin_csv(file_path)
    
    def _extract_with_langchain(self, file_path: str) -> str:
        try:
            loader = CSVLoader(file_path)
            documents = loader.load()
            return "\n\n".join([doc.page_content for doc in documents])
        except Exception as e:
            logger.error(f"Langchain CSV extraction failed: {e}")
            if PANDAS_AVAILABLE:
                return self._extract_with_pandas(file_path)
            else:
                return self._extract_with_builtin_csv(file_path)
    
    def _extract_with_pandas(self, file_path: str) -> str:
        try:
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        except Exception as e:
            logger.error(f"Pandas CSV extraction failed: {e}")
            return self._extract_with_builtin_csv(file_path)
    
    def _extract_with_builtin_csv(self, file_path: str) -> str:
        try:
            import csv
            text_lines = []
            with open(file_path, 'r', encoding='utf-8', newline='') as file:
                reader = csv.reader(file)
                for row in reader:
                    text_lines.append('\t'.join(row))
            return '\n'.join(text_lines)
        except Exception as e:
            logger.error(f"Built-in CSV extraction failed: {e}")
            raise


class XLSXExtractor(DocumentExtractor):
    """Extract text from XLSX files"""
    
    def can_handle(self, file_path: str, mime_type: str) -> bool:
        return (
            mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or
            file_path.lower().endswith('.xlsx')
        )
    
    def extract_text(self, file_path: str) -> str:
        if PANDAS_AVAILABLE:
            return self._extract_with_pandas(file_path)
        else:
            raise RuntimeError("XLSX file processing requires pandas")
    
    def _extract_with_pandas(self, file_path: str) -> str:
        try:
            # Read all sheets
            xls = pd.ExcelFile(file_path)
            all_sheets_text = []
            
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                sheet_text = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                all_sheets_text.append(sheet_text)
            
            return "\n\n".join(all_sheets_text)
        except Exception as e:
            logger.error(f"Pandas XLSX extraction failed: {e}")
            raise


class PPTXExtractor(DocumentExtractor):
    """Extract text from PPTX files"""
    
    def can_handle(self, file_path: str, mime_type: str) -> bool:
        return (
            mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or
            file_path.lower().endswith('.pptx')
        )
    
    def extract_text(self, file_path: str) -> str:
        if LANGCHAIN_AVAILABLE:
            return self._extract_with_langchain(file_path)
        elif PYTHON_PPTX_AVAILABLE:
            return self._extract_with_python_pptx(file_path)
        else:
            raise RuntimeError("No PPTX processing library available")
    
    def _extract_with_langchain(self, file_path: str) -> str:
        try:
            loader = UnstructuredPowerPointLoader(file_path)
            documents = loader.load()
            return "\n\n".join([doc.page_content for doc in documents])
        except Exception as e:
            logger.error(f"Langchain PPTX extraction failed: {e}")
            if PYTHON_PPTX_AVAILABLE:
                return self._extract_with_python_pptx(file_path)
            raise
    
    def _extract_with_python_pptx(self, file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            text_runs = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                text_runs.append(slide_text)
            
            return "\n\n".join(text_runs)
        except Exception as e:
            logger.error(f"python-pptx extraction failed: {e}")
            raise


class MP3AudioExtractor(DocumentExtractor):
    """Extract text from MP3 audio files via transcription"""
    
    def can_handle(self, file_path: str, mime_type: str) -> bool:
        return (
            mime_type == "audio/mpeg" or
            mime_type == "audio/mp3" or
            file_path.lower().endswith('.mp3')
        )
    
    def extract_text(self, file_path: str) -> str:
        try:
            # Import Google service here to avoid circular imports
            from src.services.external_apis.google_service import get_google_service
            
            # Convert MP3 to WAV for transcription if needed
            wav_path = self._prepare_audio_file(file_path)
            
            # Use Google service for transcription
            google_service = get_google_service()
            transcription_result = google_service.transcribe_audio(wav_path)
            
            # Clean up temporary file if created
            if wav_path != file_path:
                import os
                os.remove(wav_path)
            
            # Return transcript
            transcript = transcription_result.get("transcript", "")
            if not transcript:
                raise RuntimeError("No transcript generated from audio file")
            
            return transcript
            
        except Exception as e:
            logger.error(f"MP3 audio extraction failed: {e}")
            raise
    
    def _prepare_audio_file(self, file_path: str) -> str:
        """Prepare MP3 file for transcription (convert to WAV if needed)"""
        try:
            if PYDUB_AVAILABLE:
                # Convert MP3 to WAV using pydub
                audio = AudioSegment.from_mp3(file_path)
                
                # Convert to mono, 16kHz for better transcription
                audio = audio.set_channels(1)
                audio = audio.set_frame_rate(16000)
                
                # Create temporary WAV file
                import tempfile
                temp_wav = tempfile.NamedTemporaryFile(suffix='.wav', delete=False)
                audio.export(temp_wav.name, format="wav")
                
                return temp_wav.name
            else:
                # If pydub not available, try using the MP3 file directly
                # Google service will handle the conversion
                return file_path
                
        except Exception as e:
            logger.error(f"Audio preparation failed: {e}")
            # Fall back to original file
            return file_path


class PlainTextExtractor(DocumentExtractor):
    """Extract text from plain text files"""
    
    def can_handle(self, file_path: str, mime_type: str) -> bool:
        return (
            mime_type.startswith("text/") or
            file_path.lower().endswith(('.txt', '.md', '.html', '.json'))
        )
    
    def extract_text(self, file_path: str) -> str:
        encodings = ['utf-8', 'latin-1', 'ascii', 'utf-16']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding, errors='ignore') as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
            except Exception as e:
                if encoding == encodings[-1]:  # Last encoding to try
                    logger.error(f"Plain text extraction failed: {e}")
                    raise
                continue
        
        # If all encodings fail, raise error
        raise RuntimeError("Failed to read file with any supported encoding")


class DocumentTextExtractor:
    """Main document text extraction service"""
    
    def __init__(self):
        self.extractors = [
            PDFExtractor(),
            DOCXExtractor(),
            DOCExtractor(),
            CSVExtractor(),
            XLSXExtractor(),
            PPTXExtractor(),
            MP3AudioExtractor(),
            PlainTextExtractor()
        ]
        
        # Initialize text splitter for chunking
        if LANGCHAIN_AVAILABLE:
            self.text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=settings.text_chunk_size,
                chunk_overlap=settings.text_chunk_overlap,
                length_function=len
            )
        else:
            self.text_splitter = None
    
    def extract_text(self, file_path: str, mime_type: str) -> str:
        """Extract text from any supported document format"""
        try:
            # Find appropriate extractor
            for extractor in self.extractors:
                if extractor.can_handle(file_path, mime_type):
                    text = extractor.extract_text(file_path)
                    logger.info(f"Extracted {len(text)} characters from {file_path}")
                    return text
            
            # If no extractor found, raise error
            raise ValueError(f"No extractor available for file type: {mime_type}")
            
        except Exception as e:
            logger.error(f"Document text extraction failed for {file_path}: {e}")
            raise
    
    def extract_and_chunk_text(self, file_path: str, mime_type: str) -> Dict[str, Any]:
        """Extract text and split into chunks for processing"""
        try:
            # Extract full text
            full_text = self.extract_text(file_path, mime_type)
            
            # Split into chunks if langchain is available
            if self.text_splitter:
                chunks = self.text_splitter.split_text(full_text)
            else:
                # Fallback chunking
                chunk_size = settings.text_chunk_size
                chunks = [full_text[i:i+chunk_size] for i in range(0, len(full_text), chunk_size)]
            
            return {
                "full_text": full_text,
                "chunks": chunks,
                "chunk_count": len(chunks),
                "total_characters": len(full_text),
                "extraction_method": "langchain" if LANGCHAIN_AVAILABLE else "fallback"
            }
            
        except Exception as e:
            logger.error(f"Document chunking failed for {file_path}: {e}")
            raise
    
    def get_document_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract basic document metadata"""
        try:
            file_stat = os.stat(file_path)
            return {
                "file_size": file_stat.st_size,
                "file_extension": Path(file_path).suffix.lower(),
                "file_name": Path(file_path).name,
                "creation_time": file_stat.st_ctime,
                "modification_time": file_stat.st_mtime
            }
        except Exception as e:
            logger.error(f"Failed to get document metadata for {file_path}: {e}")
            return {}


# Global instance
document_extractor = DocumentTextExtractor() 